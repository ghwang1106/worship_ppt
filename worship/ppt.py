""" TODO
"""
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from worship.common import log
from worship.hymn import Hymn


class PPT:
  """ TODO
  """

  def __init__(self, layout=None):
    self.layout = layout if layout else [10, 7.5]
    self.prs = Presentation()
    self.prs.slide_width, self.prs.slide_height = Inches(
        self.layout[0]), Inches(self.layout[1])

  def make_announcement(self, announcement: list) -> None:
    for s in announcement:
      if len(s) == 0:
        pass
      elif s[0] in ['-', '*'] or re.search(r'^\S*\:', s):
        self.add_paragraph(s.replace('  ', r'\n\n'), 40, spacing=1.1)
      else:
        if s[:7] == '환영 및 심방':
          s_header, s_content = s[:7], s[8:]
        elif '  ' in s:
          s_header, s_content = s.split('  ', 1)
        else:
          s_header, s_content = s.split(' ', 1)
        self.add_textbox(s_header, [1, 1, 9, 6.5], 44)
        self.p.font.underline = True
        self.add_paragraph('', 5)
        self.add_paragraph(s_content, 40)

  def make_worship1(self, worship_order, hymn):
    log.info(worship_order)
    worship_order = [
        elem for elem in worship_order
        if elem not in ['예배기원', '인도자', '참회기도', '다함께']
    ]
    headers = ['입례송', '신앙고백', '찬양', '대표기도', '봉헌/기도']
    for s in worship_order:
      if s in headers:
        header = s
      elif header in ['입례송', '찬양']:
        if '[찬송가' in s:
          hymn_title, hymn_num = [
              x.strip().strip('"').strip('장]') for x in s.split('[찬송가')
          ]
          self.add_lyrics(hymn, title=hymn_title, number=hymn_num)
        else:
          for song_title in s.split('\n'):
            # TODO update CCM DB add code to search for CCM lyrics here
            self.add_textbox(song_title, [1, 1, 9, 6.5], 44, spacing=1.1)
        self.add_slide()
      elif header == '신앙고백':
        d = [1, 1.8, 9, 5.7]
        creed_text = [
            '신앙고백 – 사도신경',
            '전능하사 천지를 만드신\n하나님 아버지를 내가 믿사오며\n그 외아들 우리 주 예수\n그리스도를 믿사오니,',
            '이는 성령으로 잉태하사\n동정녀 마리아에게 나시고\n본디오 빌라도에게 고난을 받으사,\n십자가에 못 박혀 죽으시고,',
            '장사한지 사흘 만에\n죽은 자 가운데서 다시 살아나시며,\n하늘에 오르사 전능하신\n하나님 우편에 앉아 계시다가,',
            '저리로서 산 자와 죽은 자를\n심판하러 오시리라.\n성령을 믿사오며 거룩한 공회와,\n성도가 서로 교통하는 것과,',
            '죄를 사하여 주시는 것과,\n몸이 다시 사는 것과,\n영원히 사는 것을 믿사옵나이다.\n아멘'
        ]
        self.add_textbox(creed_text[0],
                         [sum(x) for x in zip(d, [0, -0.8, 0, 0])],
                         44,
                         align=PP_ALIGN.CENTER)
        self.add_paragraph('', 5)
        self.add_paragraph(creed_text[1],
                           40,
                           align=PP_ALIGN.CENTER,
                           spacing=1.1)
        for ac in range(2, 6):
          self.add_textbox(creed_text[ac],
                           d,
                           40,
                           align=PP_ALIGN.CENTER,
                           spacing=1.1)
        self.add_slide()
      elif header == '대표기도':
        self.add_textbox('대표 기도', [1, 1, 9, 6.5], 44, new_slide=False)
        self.add_paragraph(s, 44)
      elif header == '봉헌/기도':
        # TODO Allow lyric variation (only sing 1,3,4 lines of hymn 50 verse 1)
        self.add_slide()
        self.add_textbox('봉  헌', [1, 1, 9, 6.5], 48)
        if '[찬송가' in s:
          hymn_title, hymn_num = [
              x.strip().strip('"').strip('장]') for x in s.split('[찬송가')
          ]
          self.add_lyrics(hymn, title=hymn_title, number=hymn_num, verses=[1])
        else:  # some days we have solos
          self.add_textbox(s, [1, 1, 9, 6.5], 44, spacing=1.1)
      else:  # catch all - perhaps add an error message?
        self.add_textbox(s, [1, 1, 9, 6.5], 44, spacing=1.1)
        self.add_slide()

  def make_hymn(self, hymn, hymn_no) -> None:
    for i in enumerate(re.split(r'\D+', hymn_no.strip())):
      if 1 <= int(i[1]) <= 645:
        self.add_lyrics(hymn, number=i[1])
        self.add_slide()

  def add_lyrics(self,
                 hymn: Hymn,
                 song_type='hymn',
                 title='',
                 number=0,
                 verses='all') -> None:
    song_received = hymn.lookup(song_type=song_type,
                                title=title,
                                number=int(number))
    lyrics = song_received['lyrics']
    if verses != 'all':
      lyrics = [lyrics[v - 1] for v in verses]

    self.add_textbox(song_received['title'].split(' [')[0], [1, 0.4, 9, 1],
                     44,
                     spacing=1)
    self.tx_box.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    self.p.font.underline = True
    if song_type == 'hymn':
      self.add_paragraph('(찬송가' + str(song_received['number']) + ')',
                         28,
                         spacing=1)
      self.p.space_before = Pt(3)
    lyrics = [x.split('\n\n') for x in lyrics]
    lyrics = [item for elem in lyrics for item in elem]
    self.add_textbox(lyrics[0], [1, 1.8, 9, 3.76], 44, new_slide=False)
    for h in range(1, len(lyrics)):
      self.add_textbox(lyrics[h], [1, 1.8, 9, 3.76], 44)

  def add_slide(self, n_slide: int = 1) -> None:  # create new blank slide
    for _ in range(n_slide):
      self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
      self.slide.background.fill.solid()

  def add_textbox(self,
                  text,
                  dim,
                  font_size,
                  font_name='Malgun Gothic',
                  spacing=1.1,
                  bold=True,
                  align=PP_ALIGN.CENTER,
                  new_slide=True):
    if new_slide:
      self.add_slide()
    self.tx_box = self.slide.shapes.add_textbox(Inches(dim[0]), Inches(dim[1]),
                                                Inches(dim[2]), Inches(dim[3]))
    self.p = self.tx_box.text_frame.paragraphs[0]
    self.format_text(text, font_size, font_name, spacing, bold, align)

  def add_paragraph(self,
                    text,
                    font_size,
                    font_name='Malgun Gothic',
                    spacing=1.1,
                    bold=True,
                    align=PP_ALIGN.CENTER):
    self.p = self.tx_box.text_frame.add_paragraph()
    self.format_text(text, font_size, font_name, spacing, bold, align)

  def format_text(self,
                  text,
                  font_size,
                  font_name='Malgun Gothic',
                  spacing=1.1,
                  bold=True,
                  align=PP_ALIGN.CENTER):
    self.p.text = text
    self.p.font.name = font_name
    self.p.font.size = Pt(font_size)
    self.p.alignment = align
    self.p.space_before = Pt(6)
    self.p.line_spacing = spacing
    self.p.font.bold = bold
    self.p.font.color.rgb = RGBColor(255, 255, 255)
    self.tx_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    self.tx_box.text_frame.word_wrap = True

  def text_length(self, v):
    l = [i * 28 / 22 for i in [0.32, 0.14, 0.08]
        ]  # Malgun Gothic character lengths in inches (font_size = 28)
    n_all, n_num, n_spe = len(v), len(re.sub(r'[\D]+', '',
                                             v)), len(re.sub(r'[\w]+', '', v))
    return (n_all - n_spe - n_num) * l[0] + n_num * l[1] + n_spe * l[2]

  def to_pptx(self, path: str) -> str:
    try:
      self.prs.save(path + '.pptx')  # save to pptx with specified path
    except PermissionError:
      path = path + '1'
      self.to_pptx(path)
    return path
