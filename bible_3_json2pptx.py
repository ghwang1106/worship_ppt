import os
import re
import math
import json
import pandas as pd
from math import floor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs1 = Presentation()
prs2 = Presentation()

# Slide Size
w, h = 12192000, 6858000
prs1.slide_width = prs2.slide_width = w  # 13.33 inches
prs1.slide_height = prs2.slide_height = h  # 7.5 inches

# Slide Format


class create_slide:

    def __init__(self, prs):
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])

    def add_background(self, img_bg, date_type, passage, img_lg):
        self.bg = self.slide.shapes.add_picture(img_bg, 0, -Inches(0.66), width=w)
        self.add_textbox(date_type, [0.48, 0.19, 6.62, 0.6], 16)
        self.make_bold()
        self.p.alignment = PP_ALIGN.LEFT
        self.add_textbox(passage, [0.17, 0.64, 13, 0.68], 32, fontname='Gulim', spacing=1.5)
        self.make_bold()
        self.lg = self.slide.shapes.add_picture(img_lg, Inches(9.91), Inches(0.29), width=Inches(2.88))

    def add_textbox(self, text, dim, fontsize, fontname='Malgun Gothic', spacing=1):
        self.txBox = self.slide.shapes.add_textbox(Inches(dim[0]), Inches(dim[1]), Inches(dim[2]), Inches(dim[3]))
        self.p = self.txBox.text_frame.paragraphs[0]
        self.p.text = text
        self.p.line_spacing = spacing
        self.format_text(fontsize, fontname)

    def add_paragraph(self, text, fontsize, fontname='Malgun Gothic', spacing=1):
        self.p = self.txBox.text_frame.add_paragraph()
        self.p.text = text
        self.p.line_spacing = spacing
        self.format_text(fontsize, fontname)

    def make_bold(self):
        self.p.font.bold = True

    def add_table(self, abbr, verses):
        r, c = int(floor(len(verses)/5)+3), 5
        x, y, w, h = 1.67, 1.75, 10, 0.41*r
        self.tbl = self.slide.shapes.add_table(r, c, Inches(x), Inches(y), Inches(w), Inches(h))

        for i, v in enumerate(verses):
            rr, cc = int(floor(i/c)), i % c
            self.p = self.tbl.table.cell(rr, cc).text_frame.paragraphs[0]
            self.p.text = abbr + str(v[0]) + ":" + str(v[1])
            self.format_text(15)
            self.shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+cc*w/c), Inches(y+h*rr/r), Inches(w/c), Inches(h/r))
            self.shape.fill.background()
            self.shape.click_action.target_slide = prs1.slides[i+1]

        self.p = self.tbl.table.cell(-1, 0).text_frame.paragraphs[0]
        self.p.text = "Title"
        self.format_text(15)
        self.shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y+h*(r-1)/r), Inches(w/c), Inches(h/r))
        self.shape.fill.background()
        self.shape.click_action.target_slide = prs1.slides[0]

    def format_text(self, fontsize, fontname='Malgun Gothic'):
        self.p.font.name = fontname
        self.p.font.size = Pt(fontsize)
        self.p.alignment = PP_ALIGN.CENTER
        self.txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def parse_cv(cv):
    return [int(cv[:cv.find(":")]), int(cv[cv.find(":")+1:])]


def parse_passage(passage):
    book, cv_range = passage.split()

    if "-" in cv_range:
        tmp_from, tmp_to = cv_range.split("-")
        cv_from = parse_cv(tmp_from)
        if ":" not in tmp_to:
            tmp_to = str(cv_from[0]) + ":" + tmp_to
        cv_to = parse_cv(tmp_to)
    else:
        cv_from = parse_cv(cv_range)
        cv_to = parse_cv(cv_range)

    bb = pd.read_csv('bible_info.csv')
    i_book = int(bb[bb.Full == book].index[0])
    abbr = bb.Abbr[i_book]
    book_eng = bb.Eng[i_book]

    return cv_from, cv_to, i_book, abbr, book_eng


def read_from_json(book_eng, cv_from, cv_to):
    with open("bible.json", "r", encoding='utf-8') as f:
        dic = json.load(f)[book_eng]

    c_lookup = cv_from[0]
    v_lookup = cv_from[1]
    cv_all = []
    while c_lookup <= cv_to[0]:
        if c_lookup == cv_to[0]:
            v_lookup_to = cv_to[1]
        else:
            v_lookup_to = len(dic[c_lookup-1])

        while v_lookup <= v_lookup_to:
            cv_all.append([c_lookup, v_lookup])
            v_lookup += 1

        c_lookup += 1
        v_lookup = 1

    verses = []
    for cv in cv_all:
        if cv_from[0] == cv_to[0]:
            verses.append(str(cv[1]) + ". " + dic[cv[0]-1][cv[1]-1])
        else:
            verses.append(str(cv[0]) + ":" + str(cv[1]) + " " + dic[cv[0]-1][cv[1]-1])
    return verses, cv_all


def verse_length(v, middle=0):
    if middle:
        l = [0.32, 0.14, 0.08]
    else:
        l = [0.275, 0.15, 0.1]
    n_all = len(v)
    n_num = len(re.sub('[\D]+', '', v))
    n_spe = len(re.sub('[\w]+', '', v))
    return n_all, (n_all-n_spe-n_num)*l[0] + n_num*l[1] + n_spe*l[2]


def verse_cut(v, middle=0):
    line_length = 12
    n_all, v_length = verse_length(v, middle)

    if v_length > line_length:
        if middle:
            n_cuts = v_length // line_length
            for c in range(math.ceil(n_all/(n_cuts+1))+2, n_all-1, math.ceil(n_all/(n_cuts+1))):
                i_cut = v.find(" ", c)
                v = v[:i_cut] + "\n" + v[i_cut+1:]
        else:
            i = 0
            i_cut = 0
            while verse_length(v[i:])[1] > line_length:
                i_cut += 1
                if verse_length(v[i:i_cut])[1] > line_length:
                    v = v[:i_space] + "\n     " + v[i_space+1:]
                    i, i_cut = i_space+7, i_space+7
                elif v[i_cut] == " ":
                    i_space = i_cut
            n_cuts = v.count("\n")
    else:
        n_cuts = 0

    return v, int(n_cuts)


# Input arguments
sermon_title = "십자가 in 신앙고백"
passage = "열왕기하 6:20-33"
preacher = "이찬우 목사"
date_type = "2021년 1월 20일 수요예배"

# Verses
[cv_from, cv_to, i_book, abbr, book_eng] = parse_passage(passage)
verses, cv_all = read_from_json(book_eng, cv_from, cv_to)

# First pptx - per verse
tSlide = create_slide(prs1)
tSlide.add_textbox('"' + sermon_title + '"', [1.67, 0.08, 10, 1.2], 26, spacing=1.1)
tSlide.make_bold()
tSlide.add_paragraph('(' + passage + '/ ' + preacher + ')', 24)
tSlide.make_bold()

vSlides = []
for i, v in enumerate(verses):
    vSlides.append(create_slide(prs1))
    v, n_cuts = verse_cut(v, 1)
    if n_cuts > 2:
        v_parts = v.split("\n", n_cuts-1)
        v_half = len(v) - len(v_parts[-1]) - 2
        vSlides[-1].add_textbox(v[:v_half] + "-", [1.67, 0.08, 10, 1.2], 22, spacing=1.1)
        vSlides.append(create_slide(prs1))
        vSlides[-1].add_textbox(v.split(". ")[0] + ". " + v[v_half+2:], [1.67, 0.08, 10, 1.2], 22, spacing=1.1)
        cv_all = cv_all[:i] + [cv_all[i][:]] + cv_all[i:]
        cv_all[i][1] = str(cv_all[i][1]) + "a"
        cv_all[i+1][1] = str(cv_all[i+1][1]) + "b"
    else:
        vSlides[-1].add_textbox(v, [1.67, 0.08, 10, 1.2], 22, spacing=1.1)
tSlide.add_table(abbr, cv_all)
for i in range(len(vSlides)):
    vSlides[i].add_table(abbr, cv_all)

# Second pptx - all passage
n_lines = 7
for v in verses:
    if v[1] == ".":
        v = v[:2] + " " + v[2:]
    v, n_cuts = verse_cut(v)
    n_lines += n_cuts+1
    if n_lines > 7:
        aSlide = create_slide(prs2)
        aSlide.add_background('bible_background.jpg', date_type, passage, 'church_logo.png')
        aSlide.add_textbox(v, [0.47, 1.71, 12.4, 5.04], 20, fontname='Gulim', spacing=1.8)
        n_lines = n_cuts+1
    else:
        aSlide.add_paragraph(v, 20, fontname='Gulim', spacing=2)
    aSlide.p.alignment = PP_ALIGN.LEFT
    aSlide.p.space_before = Pt(6)
    aSlide.txBox.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    aSlide.make_bold()

# Save to File
prs1.save("Bible.pptx")
os.startfile("Bible.pptx")

date = re.split("[^0-9]", date_type)[0:5:2]
if len(date[1]) == 1:
    date[1] = '0' + date[1]
if len(date[2]) == 1:
    date[2] = '0' + date[2]
prs2.save('.'.join(date) + "(Large).pptx")
os.startfile('.'.join(date) + "(Large).pptx")
