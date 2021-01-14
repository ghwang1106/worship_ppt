import os
import pandas as pd
from math import floor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from selenium import webdriver
from selenium.webdriver.chrome.options import Options

# Bible URL
url = 'https://www.bskorea.or.kr/bible/korbibReadpage.php?version=GAE&book=num&chap=9&sec=1'

# Initialize Presentation
prs1 = Presentation()

# Slide Size
prs1.slide_width = 12192000  # 13.33 inches
prs1.slide_height = 6858000  # 7.5 inches

# Slide Format


class create_slide:

    def __init__(self, text, font):
        self.slide = prs1.slides.add_slide(prs1.slide_layouts[6])
        self.txBox = self.slide.shapes.add_textbox(Inches(1.67), Inches(0.08), Inches(10), Inches(1.2))
        self.tf = self.txBox.text_frame

        self.p = self.tf.paragraphs[0]
        self.p.text = text
        self.format_text(font)
        self.p.line_spacing = 1.1

    def make_bold(self):
        self.p.font.bold = True

    def add_paragraph(self, text, font):
        self.p = self.tf.add_paragraph()
        self.p.text = text[0]
        self.format_text(font)

    def add_table(self, passage):

        vall, chapter, abbr = parse_passage(passage)[0:3]

        prefix = abbr + chapter + ":"

        r, c = int(floor(len(vall)/5)+3), 5
        x, y, w, h = 1.67, 1.75, 10, 0.41*r
        self.tbl = self.slide.shapes.add_table(r, c, Inches(x), Inches(y), Inches(w), Inches(h))

        for i in range(len(vall)):
            rr, cc = int(floor(i/c)), i % c
            self.p = self.tbl.table.cell(rr, cc).text_frame.paragraphs[0]
            self.p.text = prefix + str(vall[i])
            self.format_text(15)
            self.shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+cc*w/c), Inches(y+h*rr/r), Inches(w/c), Inches(h/r))
            self.shape.fill.background()
            self.shape.click_action.target_slide = prs1.slides[i+1]

        self.p = self.tbl.table.cell(-1, 0).text_frame.paragraphs[0]
        self.p.text = "Title"
        self.format_text(15)
        self.shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y+h*(r-1)/r), Inches(w/c), Inches(h/r))
        self.shape.fill.background()
        self.shape.click_action.target_slide = prs1.slides[0]

    def format_text(self, font):
        self.p.font.name = 'Malgun Gothic'
        self.p.font.size = Pt(font)
        self.p.alignment = PP_ALIGN.CENTER
        self.tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def parse_passage(passage):
    book = passage[:passage.find(" ")]
    chapter = passage[passage.find(" ")+1:passage.find(":")]
    bb = pd.read_csv('bible_chapters.csv')
    abbr = bb.Abbr[bb[bb.Full == book].index].values[0]

    vrange = passage[passage.find(":")+1:]
    v1 = vrange[:vrange.find("-")]
    v2 = vrange[vrange.find("-")+1:]
    vall = range(int(v1), int(v2)+1)

    return vall, chapter, abbr, book


def extract_verses(text_all, vall):
    verses = []
    for v in vall:
        verses.append([text_all[text_all.find(str(v)):text_all.find(str(v+1))].strip().replace("   ", ". ")])

    return verses


# First Slide
sermon_title = '"목적지보다 동행자가 중요하다!"'
passage = "민수기 9:16-23"
preacher = "이신형 목사"

tSlide = create_slide(sermon_title, 26)
tSlide.make_bold()
tSlide.add_paragraph(['(' + passage + '/ ' + preacher + ')'], 24)
tSlide.make_bold()

# Verses
vall = parse_passage(passage)[0]

chrome_options = Options()
chrome_options.add_argument("--incognito")
chrome_options.add_argument("--window-size=1920x1080")
driver = webdriver.Chrome(options=chrome_options, executable_path="C:/Users/ghwan/OneDrive/Desktop/chromedriver.exe")
driver.get(url)
text_all = driver.find_element_by_class_name("bible_read").text
verses = extract_verses(text_all, vall)
driver.close()

vSlides = []
for v in verses:
    if len(v[0]) > 50:
        index_cut = v[0].find(" ", int(round(len(v[0])/2, 0))+5)
        v[0] = v[0][: index_cut] + "\n" + v[0][index_cut+1:]

    vSlides.append(create_slide(v[0], 22))

# Link Table
tSlide.add_table(passage)
for i in range(len(vSlides)):
    vSlides[i].add_table(passage)

# Save to File
prs1.save("Bible.pptx")
os.startfile("Bible.pptx")
