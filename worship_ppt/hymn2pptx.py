"""
TODO
"""
import re
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from worship_ppt.common import DATA_PATH, log


class Hymn():
  """ TODO
  """

  def __init__(self, json_file=DATA_PATH / 'hymn.json'):
    with open(json_file, 'r', encoding='utf-8') as f:
      self.all = json.load(f)

  def lookup(self, song_type='hymn', number=0, title=''):
    try:
      song_received = next(item for item in self.all
                           if item['song_type'] == song_type and
                           item['number'] == number and title in item['title'])
    except StopIteration:
      log.warning('Song not found (%s)', locals()['title'])
      song_received = {
          'song_type': song_type,
          'number': number,
          'title': title,
          'lyrics': ['(가사가 DB에 없습니다)']
      }
    return song_received


class PPT:
  """ TODO
  """

  def __init__(self, layout=None):
    self.layout = [10, 7.5] if layout is None else layout
    self.prs = Presentation()
    self.prs.slide_width, self.prs.slide_height = Inches(
        self.layout[0]), Inches(self.layout[1])

  def make_hymn(self, hymn, hymn_no):
    for i in enumerate(re.split(r'\D+', hymn_no.strip())):
      if int(i[1]) >= 1 & int(i[1]) <= 645:
        self.add_lyrics(hymn, number=i[1])
        self.add_slide()

  def add_lyrics(self,
                 hymn,
                 song_type='hymn',
                 title='',
                 number=0,
                 verses='all'):
    song_received = hymn.lookup(song_type=song_type,
                                title=title,
                                number=int(number))
    lyrics = song_received['lyrics']
    if verses == 'all':
      pass
    else:
      lyrics = [lyrics[v - 1] for v in verses]

    lines = sum([l.split('\n') for l in lyrics], [])
    w = min(max(max([self.text_length(line) for line in lines]) + 0.3, 5), 10)

    self.add_textbox(song_received['title'].split(' [')[0],
                     [10 - w, 2.65, w, 1],
                     28,
                     spacing=1)
    self.tx_box.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    self.p.font.underline = True
    if song_type == 'hymn':
      self.add_paragraph('찬송가' + str(song_received['number']), 20, spacing=1)
      self.p.space_before = Pt(3)
    lyrics = [x.split('\n\n') for x in lyrics]
    lyrics = [item for elem in lyrics for item in elem]
    self.add_textbox(lyrics[0], [10 - w, 3.7, w, 3.76], 28, newslide=0)
    for h in range(1, len(lyrics)):
      self.add_textbox(lyrics[h], [10 - w, 3.7, w, 3.76], 28)

  def add_slide(self, n_slide=1):  # create new blank slide
    for _ in range(n_slide):
      self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
      self.slide.background.fill.solid()

  def add_textbox(self,
                  text,
                  dim,
                  fontsize,
                  fontname='Malgun Gothic',
                  spacing=1.1,
                  bold=True,
                  align=PP_ALIGN.CENTER,
                  newslide=1):
    if newslide:
      self.add_slide()
    self.tx_box = self.slide.shapes.add_textbox(Inches(dim[0]), Inches(dim[1]),
                                                Inches(dim[2]), Inches(dim[3]))
    self.p = self.tx_box.text_frame.paragraphs[0]
    self.format_text(text, fontsize, fontname, spacing, bold, align)

  def add_paragraph(self,
                    text,
                    fontsize,
                    fontname='Malgun Gothic',
                    spacing=1.1,
                    bold=True,
                    align=PP_ALIGN.CENTER):
    self.p = self.tx_box.text_frame.add_paragraph()
    self.format_text(text, fontsize, fontname, spacing, bold, align)

  def format_text(self,
                  text,
                  fontsize,
                  fontname='Malgun Gothic',
                  spacing=1.1,
                  bold=True,
                  align=PP_ALIGN.CENTER):
    self.p.text = text
    self.p.font.name = fontname
    self.p.font.size = Pt(fontsize)
    self.p.alignment = align
    self.p.space_before = Pt(6)
    self.p.line_spacing = spacing
    self.p.font.bold = bold
    self.p.font.color.rgb = RGBColor(255, 255, 255)
    self.tx_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    self.tx_box.text_frame.word_wrap = True

  def make_bold(self):
    self.p.font.bold = True

  def text_length(self, v):
    l = [i * 28 / 22 for i in [0.32, 0.14, 0.08]
        ]  # Malgun Gothic character lengths in inches (fontsize = 28)
    n_all, n_num, n_spe = len(v), len(re.sub(r'[\D]+', '',
                                             v)), len(re.sub(r'[\w]+', '', v))
    return (n_all - n_spe - n_num) * l[0] + n_num * l[1] + n_spe * l[2]

  def to_pptx(self, path):
    try:
      self.prs.save(path + '.pptx')  # save to pptx with specified path
    except PermissionError:
      path = path + '1'
      self.to_pptx(path)
    return path


def make_hymn_ppt(hymn_no):
  hymn = Hymn()
  worship_ppt = PPT()
  worship_ppt.make_hymn(hymn, hymn_no)

  file_name = DATA_PATH / 'hymn'
  file_name = worship_ppt.to_pptx(file_name.as_posix())

  return file_name
