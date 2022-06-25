"""
TODO
"""
import re
import json
from itertools import groupby
from docx import Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from worship_ppt.common import DATA_PATH, log


class Hymn():
  """
  TODO
  """

  def __init__(self, json_file=DATA_PATH / 'hymn.json'):
    with open(json_file, 'r', encoding='utf-8') as f:
      self.all = json.load(f)

  def lookup(self, song_type='hymn', number=0, title=''):
    try:
      song_received = next(item for item in self.all
                           if item['song_type'] == song_type and
                           item['number'] == number and title in item['title'])
    except StopIteration:
      print('  >>WARNING: Song not found (' + locals()['title'] + ')')
      song_received = {
          'song_type': song_type,
          'number': number,
          'title': title,
          'lyrics': ['(가사가 DB에 없습니다)']
      }
    return song_received


class Jubo():
  """ TODO
  """

  def __init__(self, jubo_file):
    jubo_raw = Document(jubo_file)
    all_text = []
    for block in self.iter_block_items(jubo_raw):
      if not isinstance(block, Table):
        continue
      for row in block.rows:
        row_data = [cell.text for cell in row.cells]
        if not row_data:
          continue
        for item in set(row_data):
          all_text.append(item)
    self.all_text = [x[0] for x in groupby(all_text)]

  def iter_block_items(self, jubo_raw):
    for content in jubo_raw.element.body.iterchildren():
      if isinstance(content, CT_P):
        yield Paragraph(content, jubo_raw)
      elif isinstance(content, CT_Tbl):
        yield Table(content, jubo_raw)

  def grab_section(self, s_start, s_finish):
    if isinstance(int, s_finish):
      return self.all_text[self.find_index(self.all_text, s_start)[0]:self.
                           find_index(self.all_text, s_start)[0] + s_finish]
    else:
      return self.all_text[self.find_index(self.all_text, s_start)[0]:self.
                           find_index(self.all_text, s_finish)[0]]

  def find_index(self, l, s):
    return [i for i, ll in enumerate(l) if s in ll]


class PPT:
  """
  TODO
  """

  def __init__(self, layout=None):
    self.layout = layout if layout else [10, 7.5]
    self.prs = Presentation()
    self.prs.slide_width, self.prs.slide_height = Inches(layout[0]), Inches(
        layout[1])

  def make_announcement(self, announcement):
    for s in announcement:
      if len(s) == 0:
        pass
      elif s[0] in ['-', '*'] or re.search(r'^\S*\:', s):
        self.add_paragraph(s.replace('  ', r'\n\n'), 28, spacing=1.1)
      else:
        if '  ' in s:
          print(s)
          s_header, s_content = s.split('  ', 1)
        else:
          s_header, s_content = s.split(' ', 1)
        self.add_textbox(s_header, [2.5, 2.51, 7.5, 4.95], 32)
        self.add_paragraph('', 5)
        self.add_paragraph(s_content, 28)

  def make_worship1(self, worship_order, hymn):
    f_empty = 0
    f_cong_pray = 0
    f_offering = 0
    for s in worship_order:
      to_ignore = ['입례송', '예배기원', '인도자', '신앙고백', '찬양', '참회기도', '다함께']
      if s in to_ignore:
        if not f_empty:
          f_empty = 1
          self.add_slide()
      elif '[찬송가' in s:
        hymn_title, hymn_num = [
            x.strip().strip('"').strip('장]') for x in s.split('[찬송가')
        ]
        self.add_lyrics(hymn, title=hymn_title, number=hymn_num)
        self.add_slide()
      elif s == '사도신경':
        d = [3.09, 3.7, 6.8, 3.76]
        creed_text = [
            '신앙고백-사도신경',
            '전능하사 천지를 만드신\n하나님 아버지를 내가 믿사오며,\n그 외아들 우리 주 예수\n그리스도를 믿사오니,',
            '이는 성령으로 잉태하사\n동정녀 마리아에게 나시고,\n본디오 빌라도에게 고난을 받으사,\n십자가에 못 박혀 죽으시고',
            '장사한지 사흘 만에\n죽은 자 가운데서 다시 살아나시며,\n하늘에 오르사 전능하신 하나님\n우편에 앉아 계시다가,',
            '저리로서 산 자와 죽은 자를\n심판하러 오시리라.\n성령을 믿사오며 거룩한 공회와,\n성도가 서로 교통하는 것과,',
            '죄를 사하여 주시는 것과,\n몸이 다시 사는 것과,\n영원히 사는 것을 믿사옵나이다.\n아멘'
        ]
        self.add_textbox(creed_text[0],
                         [sum(x) for x in zip(d, [0, -0.84, 0, 0])], 32)
        self.add_paragraph('', 5)
        self.add_paragraph(creed_text[1], 28)
        for ac in range(2, 6):
          self.add_textbox(creed_text[ac], d, 28)
        self.add_slide()
      elif s == '대표기도':
        f_cong_pray = 1
        self.add_textbox('대표 기도\n', [0, 0.95, 10, 2.25], 44, newslide=0)
      elif f_cong_pray:
        f_cong_pray = 0
        self.add_paragraph(s, 44)
        self.add_slide()
      elif s == '봉헌/기도':
        f_offering = 1
      elif f_offering:
        f_offering = 0
        self.add_lyrics(hymn, song_type='ccm', title='나의 모습 나의 소유', verses=[4])
        self.add_slide()
      else:
        f_empty = 0
        self.add_textbox(s, [3.09, 2.51, 6.93, 4.95], 28, spacing=1.1)
        self.add_slide()

  def make_worship2(self, title):

    title = title.split('\n')
    if len(title) != 2:
      log.warning('Unexpected number of lines in the title page')

    self.add_textbox('\n'.join(title[0:-1]), [0, 1.5, 10, 1.5], 36)
    self.add_paragraph(title[-1], 32)
    self.add_slide()
    self.add_textbox('축  도', [0, 1.3, 10, 2.25], 44)

  def add_lyrics(self,
                 hymn,
                 song_type='hymn',
                 title='',
                 number=0,
                 verses='all'):
    song_received = hymn.lookup(song_type=song_type,
                                title=title,
                                number=int(number))
    lyrics = song_received['lyrics']
    if verses == 'all':
      pass
    else:
      lyrics = [lyrics[v - 1] for v in verses]

    lines = sum([l.split('\n') for l in lyrics], [])
    w = min(max(max([self.text_length(line) for line in lines]) + 0.3, 5), 10)

    self.add_textbox(song_received['title'].split(' [')[0],
                     [10 - w, 2.65, w, 1],
                     28,
                     spacing=1)
    self.tx_box.text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    self.p.font.underline = True
    if song_type == 'hymn':
      self.add_paragraph('찬송가' + str(song_received['number']), 20, spacing=1)
      self.p.space_before = Pt(3)
    lyrics = [x.split('\n\n') for x in lyrics]
    lyrics = [item for elem in lyrics for item in elem]
    self.add_textbox(lyrics[0], [10 - w, 3.7, w, 3.76], 28, newslide=0)
    for h in range(1, len(lyrics)):
      self.add_textbox(lyrics[h], [10 - w, 3.7, w, 3.76], 28)

  def add_slide(self, n=1):  # create new blank slide
    for _ in range(n):
      self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
      self.slide.background.fill.solid()

  def add_textbox(self,
                  text,
                  dim,
                  fontsize,
                  fontname='Malgun Gothic',
                  spacing=1.1,
                  bold=True,
                  align=PP_ALIGN.CENTER,
                  newslide=1):
    if newslide:
      self.add_slide()
    self.tx_box = self.slide.shapes.add_textbox(Inches(dim[0]), Inches(dim[1]),
                                                Inches(dim[2]), Inches(dim[3]))
    self.p = self.tx_box.text_frame.paragraphs[0]
    self.format_text(text, fontsize, fontname, spacing, bold, align)

  def add_paragraph(self,
                    text,
                    fontsize,
                    fontname='Malgun Gothic',
                    spacing=1.1,
                    bold=True,
                    align=PP_ALIGN.CENTER):
    self.p = self.tx_box.text_frame.add_paragraph()
    self.format_text(text, fontsize, fontname, spacing, bold, align)

  def format_text(self,
                  text,
                  fontsize,
                  fontname='Malgun Gothic',
                  spacing=1.1,
                  bold=True,
                  align=PP_ALIGN.CENTER):
    self.p.text = text
    self.p.font.name = fontname
    self.p.font.size = Pt(fontsize)
    self.p.alignment = align
    self.p.space_before = Pt(6)
    self.p.line_spacing = spacing
    self.p.font.bold = bold
    self.tx_box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    self.tx_box.text_frame.word_wrap = True
    self.p.font.color.rgb = RGBColor(255, 255, 255)

  def make_bold(self):
    self.p.font.bold = True

  def text_length(self, v):
    l = [i * 28 / 22 for i in [0.32, 0.14, 0.08]
        ]  # Malgun Gothic character lengths in inches (fontsize = 28)
    n_all, n_num, n_spe = len(v), len(re.sub(r'[\D]+', '',
                                             v)), len(re.sub(r'[\w]+', '', v))
    return (n_all - n_spe - n_num) * l[0] + n_num * l[1] + n_spe * l[2]

  def to_pptx(self, path):
    self.prs.save(path + '.pptx')
    return path  # save to pptx with specified path


def make_main_ppt(jubo_path):
  hymn = Hymn()
  jubo = Jubo(jubo_path)
  jubo.sermon = jubo.grab_section('2021년 표어', 2)[1].strip('\n')
  jubo.announcement = jubo.grab_section('우리 교회에 처음 오신 성도님들을', 3)[2].split('\n')
  jubo.worship_order = jubo.grab_section('입례송', '성경봉독')

  worship_ppt = PPT()
  worship_ppt.add_slide()
  worship_ppt.make_announcement(jubo.announcement)
  worship_ppt.make_worship1(jubo.worship_order, hymn)
  worship_ppt.make_worship2(jubo.sermon)

  file_name = DATA_PATH / 'worship'
  file_name = worship_ppt.to_pptx(file_name.as_posix())

  return file_name
