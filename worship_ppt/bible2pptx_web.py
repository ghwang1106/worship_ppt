import os
import re
import sys
import json
import pandas as pd

from copy import deepcopy
from math import floor
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from .common import DATA_PATH, log


class Bible:
  def __init__(self, bible: pd.DataFrame, json_file: Path = DATA_PATH / 'bible.json'):
    self.bible = bible
    with open(json_file, 'r', encoding='utf-8') as f:
      self.all_text = json.load(f)

  def lookup(self, verses_tf):
    verses_itr = []

    for verse_tf in verses_tf:
      if verse_tf[1][0] == verse_tf[2][0]:  # if verses are within the same chapter
        assert verse_tf[1][1] <= verse_tf[2][1], "Invalid Bible Verse : `TO` must be greater than `FROM`"
        if verse_tf[1][1] == verse_tf[2][1]:  # if only one verse
          verses_itr.append(verse_tf[:2])
        else:
          for j in range(verse_tf[2][1] - verse_tf[1][1] + 1):
            verses_itr.append([verse_tf[0], [verse_tf[1][0], verse_tf[1][1] + j]])
      else:  # if verses run across chapters
        assert verse_tf[1][0] <= verse_tf[2][0], "Invalid Bible chapter range : `TO` must be greater than `FROM`"
        c, v = verse_tf[1]
        while c <= verse_tf[2][0]:
          if c == verse_tf[2][0]:
            vt = verse_tf[2][1]
          else:
            vt = len(self.all_text[verse_tf[0]][c - 1])

          for j in range(vt - v + 1):
            verses_itr.append([verse_tf[0], [c, v + j]])
          c += 1
          v = 1
    if len(set(map(tuple, [[a[0], *a[1]] for a in verses_itr]))) < len(verses_itr):
      log.warning('There are duplicate verses')

    for i in range(len(verses_itr)):
      verses_itr[i].append(self.get_verse(verses_itr[i]))  # look up verse
      verses_itr[i] = [self.bible.Abbr[int(self.bible[self.bible.Eng == verses_itr[i][0]].index[0])]] + verses_itr[i]

    return verses_itr

  def get_verse(self, verse):
    try:
      return self.all_text[verse[0]][verse[1][0] - 1][verse[1][1] - 1]
    except IndexError:
      sys.exit('  >> ERROR: Bible verse out of index!')


class Sermon:
  def __init__(self, bible: pd.DataFrame, raw_inputs):
    self.bible = bible

    if len(raw_inputs) > 4:
      if raw_inputs[4].isspace() or raw_inputs[4] == '':  # Remove extra lines
        raw_inputs = raw_inputs[:4]

    self.date_type, self.title, self.preacher, passage = raw_inputs[:4]  # assign to variables (order must match!)
    try:
      if int(self.date_type[:3]) != 202 or sum(x.isdigit() for x in self.date_type) < 6:
        log.warnning('Check dates')
    except ValueError:
      log.warning('Check dates')
    passages_raw = [x.strip() for x in passage.split(',')]
    self.passages_raw, self.passages_ind = self.passages_raw2ind(passages_raw)

    if len(self.passages_raw) > 1:  # if there are multiple blocks of text within the main passage
      if len(list(set([a.split(' ')[0] for a in self.passages_raw]))) == 1:
        if len(list(set([a.split(' ')[-1].split(':')[0] for a in self.passages_raw]))) == 1:
          self.passages_raw[1:] = [a.split(':')[-1] for a in self.passages_raw[1:]]
        else:
          self.passages_raw[1:] = [a.split(' ')[-1] for a in self.passages_raw[1:]]
      self.passages_raw = [', '.join(self.passages_raw)]

    if len(raw_inputs) == 5:  # if there are quotes to include other than the main passage
      quotes_raw = [x.strip() for x in raw_inputs[4].split(',')]
      self.quotes_raw, self.quotes_ind = self.passages_raw2ind(quotes_raw)
    else:
      self.quotes_ind = []

  def passages_raw2ind(self, passages_raw):
    passages_ind = []

    for i, passage_raw in enumerate(passages_raw):
      assert passage_raw.count(" ") == 1, "Passage Format Error: Each passage should have one blank"
      [book, verses] = passages_raw[i].split()
      if book in list(self.bible.Abbr):  # if only abbreviation of the book is given, turn it into full name
        passages_raw[i] = self.bible.Full[int(self.bible[self.bible.Abbr == book].index[0])] + ' ' + verses
      passages_ind.append(self.parse_passage(passages_raw[i]))

    return passages_raw, passages_ind

  def parse_passage(self, passage):
    try:
      book, verse_range = passage.split()
      book_ind = self.bible.Eng[int(self.bible[self.bible.Full == book].index[0])]
    except (ValueError, IndexError) as e:
      raise f"Bible book not recognizable ({passage})" from e

    try:
      if '-' in verse_range:
        verses = verse_range.split('-')
        if ':' not in verses[1]:
          verses[1] = verses[0].split(':')[0] + ':' + verses[1]
        verses[0] = [int(i) for i in verses[0].split(':')]
        verses[1] = [int(i) for i in verses[1].split(':')]
      else:
        verses = [[int(i) for i in verse_range.split(':')] for j in range(2)]
    except (ValueError, IndexError) as e:
      raise "Bible verse not recognizable (check 4th & 5th lines in the prep file for typos)" from e

    return [book_ind, *verses]


class PPT:
  def __init__(self, layout=None):  # 13.33, 7.5 inches
    if layout is None:
      self.layout = [12192000, 6858000]
    else:
      self.layout = layout
    if os.path.exists(DATA_PATH / 'template.pptx'):
      self.prs = Presentation(DATA_PATH /
                              'template.pptx')  # template file with "browsed by individual window" turned on
      self.prs.part.drop_rel(self.prs.slides._sldIdLst[0].rId)
      del self.prs.slides._sldIdLst[0]
    else:
      self.prs = Presentation()  # starting from default settings
    self.prs.slide_width, self.prs.slide_height = self.layout

  def create_verse(self, sermon):  # Bible.pptx format (verses and hyperlinks)
    self.add_slide()  # create title slide
    if '"' in sermon.title:
      self.add_textbox(sermon.title, [1.67, 0.08, 10, 1.2], 26, spacing=1.1)
    else:
      self.add_textbox('"' + sermon.title + '"', [1.67, 0.08, 10, 1.2], 26, spacing=1.1)
    self.add_paragraph('(' + sermon.passages_raw[0] + '/ ' + sermon.preacher + ')', 24)

    p = self.add_verse_slides(sermon.passages_ext, 'p')  # add verses in main passage
    q = self.add_verse_slides(sermon.quotes_ext, 'q')  # add verses to quote

    for slide in self.prs.slides:
      self.slide = slide
      self.add_link_table(p + q)  # add hyperlink table

  def create_large(self, sermon):  # Large.pptx format
    self.add_large_slides(sermon, sermon.passages_ext, 7)

  def add_slide(self):  # create new blank slide
    self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

  def add_background(self, img_bg, date_type, passage, img_lg):  # create background with church logo
    self.bg = self.slide.shapes.add_picture(img_bg.as_posix(), 0, -Inches(0.66), width=12192000)
    self.add_textbox(date_type, [0.48, 0.19, 6.62, 0.6], 16)
    self.p.alignment = PP_ALIGN.LEFT
    self.add_textbox(passage, [0.17, 0.64, 13, 0.68], 32, fontname='Gulim', spacing=1.5)
    self.lg = self.slide.shapes.add_picture(img_lg.as_posix(), Inches(9.91), Inches(0.29), width=Inches(2.88))

  def add_textbox(self, text, dim, fontsize, fontname='Malgun Gothic', spacing=1, bold=True):
    self.txBox = self.slide.shapes.add_textbox(Inches(dim[0]), Inches(dim[1]), Inches(dim[2]), Inches(dim[3]))
    self.p = self.txBox.text_frame.paragraphs[0]
    self.format_text(text, fontsize, fontname, spacing, bold)

  def add_paragraph(self, text, fontsize, fontname='Malgun Gothic', spacing=1, bold=True):
    self.p = self.txBox.text_frame.add_paragraph()
    self.format_text(text, fontsize, fontname, spacing, bold)

  def format_text(self, text, fontsize, fontname='Malgun Gothic', spacing=1, bold=True):
    self.p.text = text
    self.p.font.name = fontname
    self.p.font.size = Pt(fontsize)
    self.p.alignment = PP_ALIGN.CENTER
    self.p.space_before = Pt(6)
    self.p.line_spacing = spacing
    self.p.font.bold = bold
    self.txBox.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

  def add_verse_slides(self, passage, v_type):
    cnt = 0
    for i in range(len(passage)):
      self.add_slide()
      if v_type == 'p':
        txt = self.verse_style(passage[i + cnt], passage)
      elif v_type == 'q':
        txt = self.verse_style(passage[i + cnt], 'book')
      txt, n_cuts = self.verse_cut(txt, 1)

      if n_cuts > 2:  # if more than 3 lines, break into two pages (no verse exceeds 6 lines)
        passage_a = deepcopy(passage[:i + cnt + 1])
        passage_a[-1][2][1] = str(passage_a[-1][2][1]) + 'a'
        passage_b = deepcopy(passage[i + cnt:])
        passage_b[0][2][1] = str(passage_b[0][2][1]) + 'b'
        passage = passage_a + passage_b
        cnt += 1

        v_parts = txt.split("\n", n_cuts - 1)
        v_half = len(txt) - len(v_parts[-1]) - 2
        self.add_textbox(txt[:v_half + 1] + "-", [1.67, 0.08, 10, 1.2], 22, spacing=1.1, bold=False)
        self.add_slide()
        self.add_textbox(txt.split(" ")[0] + " " + txt[v_half + 2:], [1.67, 0.08, 10, 1.2],
                         22,
                         spacing=1.1,
                         bold=False)
      else:
        self.add_textbox(txt, [1.67, 0.08, 10, 1.2], 22, spacing=1.1, bold=False)

    return passage

  def add_large_slides(self, sermon, passage, n_lines):
    lc = n_lines  # count number of lines
    for v in sermon.passages_ext:
      txt = self.verse_style(v, passage)
      if v[2][1] < 10:  # for verses 1-9, add a space for better alignment
        txt = txt.replace(' ', '  ', 1)
      txt, n_cuts = self.verse_cut(txt)
      lc += n_cuts + 1
      if lc > n_lines:  # create new page for every 7 lines
        self.add_slide()
        self.add_background(DATA_PATH / 'bible_background.jpg', sermon.date_type, sermon.passages_raw[0],
                            DATA_PATH / 'church_logo.png')
        self.add_textbox(txt, [0.47, 1.71, 12.4, 5.04], 20, fontname='Gulim', spacing=1.8)
        lc = n_cuts + 1
      else:
        self.add_paragraph(txt, 20, fontname='Gulim', spacing=1.8)
      self.p.alignment = PP_ALIGN.LEFT
      self.txBox.text_frame.vertical_anchor = MSO_ANCHOR.TOP

  def add_link_table(self, verses):
    r, c = int(floor((len(verses) - 1) / 5) + 3), 5  # link table with 5 columns
    x, y, w, h = 1.67, 1.75, 10, 0.41 * r
    self.tbl = self.slide.shapes.add_table(r, c, Inches(x), Inches(y), Inches(w), Inches(h))

    for i, v in enumerate(verses):
      rr, cc = int(floor(i / c)), i % c
      self.p = self.tbl.table.cell(rr, cc).text_frame.paragraphs[0]
      self.format_text(v[0] + str(v[2][0]) + ":" + str(v[2][1]), 15, bold=False)
      self.add_link_block([x + cc * w / c, y + h * rr / r, w / c, h / r], i + 1)

    self.p = self.tbl.table.cell(-1, 0).text_frame.paragraphs[0]
    self.format_text("Title", 15, bold=False)
    self.add_link_block([x, y + h * (r - 1) / r, w / c, h / r], 0)

  def add_link_block(self, pd, slide_num):
    self.shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(pd[0]), Inches(pd[1]), Inches(pd[2]),
                                             Inches(pd[3]))
    self.shape.fill.background()
    self.shape.click_action.target_slide = self.prs.slides[slide_num]  # add hyperlink to slide

  def verse_style(self, v, style):
    if style == 'book' or len(list(set([row[0] for row in style]))) > 1:  # if multiple books
      txt = v[0] + str(v[2][0]) + ':' + str(v[2][1]) + ' ' + v[3]
    elif len(list(set([row[2][0] for row in style]))) > 1:  # if multiple chapters
      txt = str(v[2][0]) + ':' + str(v[2][1]) + ' ' + v[3]
    else:
      txt = str(v[2][1]) + '. ' + v[3]
    return txt

  def verse_length(self, v, middle=0) -> int:
    if middle:
      left, mid, right = [0.32, 0.14, 0.08]  # Malgun Gothic character lengths in inches
    else:
      left, mid, right = [0.27, 0.145, 0.09]  # Gulim character lengths in inches
    n_all, n_num, n_spe = len(v), len(re.sub(r'[\D]+', '', v)), len(re.sub(r'[\w]+', '', v))
    return n_all, (n_all - n_spe - n_num) * left + n_num * mid + n_spe * right

  def verse_cut(self, v, middle=0):
    line_length = 12
    n_all, v_length = self.verse_length(v, middle)

    if v_length > line_length:  # if verse is longer than 12 inches, add breaks
      if middle:  # For verse PPT, for Gulim
        n_cuts = int((v_length + 3) // line_length)
        i_cut = v.rindex(' ')
        i = n_all
        for c in range(n_cuts):
          avg_length = self.verse_length(v[:i], 1)[1] / (n_cuts + 1 - c)
          while self.verse_length(v[v.rindex(' ', 0, i_cut) + 1:i], 1)[1] < avg_length:
            i_cut -= 1
          i = i_cut + 1
          v = v[:v.rindex(' ', 0, i)] + '\n' + v[v.rindex(' ', 0, i) + 1:]
      else:  # For large PPT, for Malgun Gothic
        i, i_cut, i_space = 0, 0, 0
        while self.verse_length(v[i:])[1] > line_length:
          i_cut += 1
          if self.verse_length(v[i:i_cut])[1] > line_length:
            v = v[:i_space] + "\n     " + v[i_space + 1:]
            i, i_cut = i_space, i_space + 7
          elif v[i_cut] == " ":
            i_space = i_cut
        n_cuts = v.count("\n")
    else:
      n_cuts = 0

    return v, int(n_cuts)

  def to_pptx(self, path: str) -> str:
    try:
      self.prs.save(path + '.pptx')  # save to pptx with specified path
    except PermissionError:
      path = path + '1'
      self.to_pptx(path)
    return path


def make_verse_ppt(ppt_inputs: list, v=0) -> str:
  bb = pd.read_csv(DATA_PATH / 'bible_info.csv')
  bible = Bible(bb)
  sermon = Sermon(bb, ppt_inputs)
  sermon.passages_ext = bible.lookup(sermon.passages_ind)
  sermon.quotes_ext = bible.lookup(sermon.quotes_ind)

  if v:
    ppt_verse = PPT()
    ppt_verse.create_verse(sermon)
    file_name = DATA_PATH / (re.sub(r"\D+", "_", sermon.date_type) + '자막')
    file_name = ppt_verse.to_pptx(file_name.as_posix())
  else:
    ppt_large = PPT()
    ppt_large.create_large(sermon)
    file_name = DATA_PATH / (re.sub(r"\D+", "_", sermon.date_type) + '본문')
    file_name = ppt_large.to_pptx(file_name.as_posix())
  return file_name
